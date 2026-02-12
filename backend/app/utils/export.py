"""Export utilities for PDF, Excel, PowerPoint"""
from typing import List, Dict, Any
from io import BytesIO
import json


async def export_to_pdf(data: Dict[str, Any], template: str = "dashboard") -> BytesIO:
    """Export data to PDF"""
    # TODO: Implement PDF export using reportlab
    # For now, return a stub
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas
    
    buffer = BytesIO()
    p = canvas.Canvas(buffer, pagesize=letter)
    
    # Add content based on template
    if template == "dashboard":
        p.drawString(100, 750, "PipelinePro Dashboard Report")
        p.drawString(100, 730, f"Generated: {data.get('generated_at', 'N/A')}")
        # Add more content based on data
    
    p.showPage()
    p.save()
    buffer.seek(0)
    return buffer


async def export_to_excel(data: List[Dict[str, Any]], filename: str = "export") -> BytesIO:
    """Export data to Excel"""
    from openpyxl import Workbook
    from openpyxl.styles import Font, Alignment
    
    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    
    # Add headers if data exists
    if data:
        headers = list(data[0].keys())
        ws.append(headers)
        
        # Style headers
        for cell in ws[1]:
            cell.font = Font(bold=True)
            cell.alignment = Alignment(horizontal="center")
        
        # Add data rows
        for row in data:
            ws.append([row.get(header, "") for header in headers])
    
    buffer = BytesIO()
    wb.save(buffer)
    buffer.seek(0)
    return buffer


async def export_to_powerpoint(data: Dict[str, Any], template: str = "dashboard") -> BytesIO:
    """Export data to PowerPoint"""
    from pptx import Presentation
    from pptx.util import Inches
    
    prs = Presentation()
    
    # Add title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "PipelinePro Report"
    subtitle.text = f"Generated: {data.get('generated_at', 'N/A')}"
    
    # Add content slides based on data
    # TODO: Add more slides based on template and data
    
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer

